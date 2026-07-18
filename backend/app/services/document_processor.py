"""
Document processing service.

Handles text extraction from various file formats (PDF, DOCX, PPTX, Excel,
CSV, and images), then chunks, embeds, and stores the content in the
vector store.
"""

import os
import uuid
import logging
from datetime import datetime, timezone
from typing import Optional

from app.config import get_settings
from app.services import embedding_service, vector_store
from app.utils.text_processing import chunk_text, clean_text

logger = logging.getLogger(__name__)

# Supported file extensions mapped to extractor functions
SUPPORTED_EXTENSIONS = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "docx",
    ".pptx": "pptx",
    ".ppt": "pptx",
    ".xlsx": "excel",
    ".xls": "excel",
    ".csv": "csv",
    ".png": "image",
    ".jpg": "image",
    ".jpeg": "image",
    ".bmp": "image",
    ".tiff": "image",
    ".tif": "image",
    ".webp": "image",
}


def detect_file_type(file_path: str) -> str:
    """
    Detect the file type from its extension.

    Args:
        file_path: Path to the file.

    Returns:
        File type string (e.g., 'pdf', 'docx', 'image').

    Raises:
        ValueError: If the file extension is not supported.
    """
    ext = os.path.splitext(file_path)[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {ext}")
    return SUPPORTED_EXTENSIONS[ext]


def extract_text_from_pdf(file_path: str) -> str:
    """
    Extract text from a PDF file using pypdf.

    Args:
        file_path: Path to the PDF file.

    Returns:
        Extracted text content.
    """
    from pypdf import PdfReader

    try:
        reader = PdfReader(file_path)
        text_parts = []
        for page_num, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
            else:
                logger.debug("Page %d has no extractable text", page_num + 1)

        text = "\n\n".join(text_parts)
        logger.info("Extracted %d characters from PDF (%d pages)", len(text), len(reader.pages))
        return text

    except Exception as e:
        logger.error("PDF extraction failed for %s: %s", file_path, str(e))
        raise RuntimeError(f"Failed to extract text from PDF: {e}") from e


def extract_text_from_docx(file_path: str) -> str:
    """
    Extract text from a DOCX file using python-docx.

    Args:
        file_path: Path to the DOCX file.

    Returns:
        Extracted text content.
    """
    from docx import Document

    try:
        doc = Document(file_path)
        text_parts = []

        # Extract from paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # Extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)

        text = "\n\n".join(text_parts)
        logger.info("Extracted %d characters from DOCX", len(text))
        return text

    except Exception as e:
        logger.error("DOCX extraction failed for %s: %s", file_path, str(e))
        raise RuntimeError(f"Failed to extract text from DOCX: {e}") from e


def extract_text_from_pptx(file_path: str) -> str:
    """
    Extract text from a PPTX file using python-pptx.

    Args:
        file_path: Path to the PPTX file.

    Returns:
        Extracted text content.
    """
    from pptx import Presentation

    try:
        prs = Presentation(file_path)
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            slide_texts.append(para_text)
                # Extract text from tables in slides
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_texts.append(row_text)

            if slide_texts:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))

        text = "\n\n".join(text_parts)
        logger.info("Extracted %d characters from PPTX (%d slides)", len(text), len(prs.slides))
        return text

    except Exception as e:
        logger.error("PPTX extraction failed for %s: %s", file_path, str(e))
        raise RuntimeError(f"Failed to extract text from PPTX: {e}") from e


def extract_text_from_excel(file_path: str) -> str:
    """
    Extract text from an Excel file using pandas.

    Args:
        file_path: Path to the Excel file.

    Returns:
        Extracted text content (all sheets concatenated).
    """
    import pandas as pd

    try:
        excel_file = pd.ExcelFile(file_path)
        text_parts = []

        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            if not df.empty:
                sheet_text = f"[Sheet: {sheet_name}]\n"
                # Include column headers
                sheet_text += " | ".join(str(col) for col in df.columns) + "\n"
                # Include rows
                for _, row in df.iterrows():
                    row_text = " | ".join(str(val) for val in row.values if pd.notna(val))
                    if row_text.strip():
                        sheet_text += row_text + "\n"
                text_parts.append(sheet_text)

        text = "\n\n".join(text_parts)
        logger.info("Extracted %d characters from Excel (%d sheets)", len(text), len(excel_file.sheet_names))
        return text

    except Exception as e:
        logger.error("Excel extraction failed for %s: %s", file_path, str(e))
        raise RuntimeError(f"Failed to extract text from Excel: {e}") from e


def extract_text_from_csv(file_path: str) -> str:
    """
    Extract text from a CSV file using pandas.

    Args:
        file_path: Path to the CSV file.

    Returns:
        Extracted text content.
    """
    import pandas as pd

    try:
        df = pd.read_csv(file_path)
        text_parts = []

        # Include column headers
        text_parts.append(" | ".join(str(col) for col in df.columns))

        # Include rows
        for _, row in df.iterrows():
            row_text = " | ".join(str(val) for val in row.values if pd.notna(val))
            if row_text.strip():
                text_parts.append(row_text)

        text = "\n".join(text_parts)
        logger.info("Extracted %d characters from CSV (%d rows)", len(text), len(df))
        return text

    except Exception as e:
        logger.error("CSV extraction failed for %s: %s", file_path, str(e))
        raise RuntimeError(f"Failed to extract text from CSV: {e}") from e


def extract_text_from_image(file_path: str) -> str:
    """
    Extract text from an image file using EasyOCR.

    Args:
        file_path: Path to the image file.

    Returns:
        Extracted text content.
    """
    from app.utils.ocr import extract_text_from_image as ocr_extract

    return ocr_extract(file_path, languages=["en"])


def _extract_text(file_path: str, file_type: str) -> str:
    """
    Route to the appropriate text extractor based on file type.

    Args:
        file_path: Path to the file.
        file_type: The detected file type string.

    Returns:
        Extracted text content.
    """
    extractors = {
        "pdf": extract_text_from_pdf,
        "docx": extract_text_from_docx,
        "pptx": extract_text_from_pptx,
        "excel": extract_text_from_excel,
        "csv": extract_text_from_csv,
        "image": extract_text_from_image,
    }

    extractor = extractors.get(file_type)
    if extractor is None:
        raise ValueError(f"No extractor available for file type: {file_type}")

    return extractor(file_path)


def process_document(
    file_path: str,
    metadata: Optional[dict] = None,
) -> dict:
    """
    Process a document end-to-end: extract text, chunk, embed, and store.

    This is the main orchestration function for document ingestion.

    Args:
        file_path: Absolute path to the uploaded file.
        metadata: Optional metadata dict (category, title, description, etc.).

    Returns:
        A dict with document_id, status, message, and chunks_created.
    """
    settings = get_settings()
    metadata = metadata or {}
    filename = os.path.basename(file_path)
    document_id = metadata.get("document_id", str(uuid.uuid4()))

    try:
        # 1. Detect file type
        file_type = detect_file_type(file_path)
        logger.info("Processing document: %s (type: %s)", filename, file_type)

        # 2. Extract text
        raw_text = _extract_text(file_path, file_type)
        if not raw_text or not raw_text.strip():
            return {
                "document_id": document_id,
                "status": "failed",
                "message": "No text could be extracted from the document",
                "chunks_created": 0,
            }

        # 3. Clean text
        cleaned_text = clean_text(raw_text)

        # 4. Chunk text
        chunks = chunk_text(
            cleaned_text,
            chunk_size=settings.chunk_size,
            chunk_overlap=settings.chunk_overlap,
        )

        if not chunks:
            return {
                "document_id": document_id,
                "status": "failed",
                "message": "Text was extracted but could not be chunked",
                "chunks_created": 0,
            }

        # 5. Generate embeddings
        embeddings = embedding_service.embed_batch(chunks)

        # 6. Build metadata for each chunk
        now = datetime.now(timezone.utc).isoformat()
        chunk_metadatas = []
        chunk_ids = []
        for i in range(len(chunks)):
            chunk_meta = {
                "source": filename,
                "source_type": file_type,
                "document_id": document_id,
                "chunk_index": i,
                "total_chunks": len(chunks),
                "upload_date": now,
                "last_updated": now,
                "category": metadata.get("category", "general"),
                "title": metadata.get("title", filename),
                "description": metadata.get("description", ""),
            }
            chunk_metadatas.append(chunk_meta)
            chunk_ids.append(f"{document_id}_chunk_{i}")

        # 7. Store in vector store
        vector_store.add_documents(
            texts=chunks,
            metadatas=chunk_metadatas,
            ids=chunk_ids,
            embeddings=embeddings,
        )

        logger.info(
            "Document '%s' processed: %d chunks created and stored",
            filename,
            len(chunks),
        )

        return {
            "document_id": document_id,
            "status": "completed",
            "message": f"Successfully processed {filename}: {len(chunks)} chunks created",
            "chunks_created": len(chunks),
        }

    except Exception as e:
        logger.error("Document processing failed for %s: %s", filename, str(e))
        return {
            "document_id": document_id,
            "status": "failed",
            "message": f"Processing failed: {str(e)}",
            "chunks_created": 0,
        }
